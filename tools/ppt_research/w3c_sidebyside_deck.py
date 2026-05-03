#!/usr/bin/env python3
"""Build a side-by-side W3C corpus deck: oracle/PPT reference vs svg2ooxml render."""

from __future__ import annotations

import argparse
import logging
from dataclasses import dataclass
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
except ImportError as exc:  # pragma: no cover - dependency check
    Presentation = None  # type: ignore[assignment]
    RGBColor = None  # type: ignore[assignment]
    PP_ALIGN = None  # type: ignore[assignment]
    Inches = None  # type: ignore[assignment]
    Pt = None  # type: ignore[assignment]
    _PPTX_IMPORT_ERROR = exc
else:
    _PPTX_IMPORT_ERROR = None

from tools.visual.builder import PptxBuilder
from tools.pptx_builder import embed as embed_svg_collection
from tools.visual.corpus_audit import (
    _stale_w3c_reference_reason,
    _w3c_reference_png_for_svg,
)
from tools.visual.renderer import PptxRenderer, resolve_renderer

logger = logging.getLogger(__name__)

BASE_SLIDE_WIDTH_IN = 13.333
SLIDE_WIDTH_IN = BASE_SLIDE_WIDTH_IN * 2.0
SLIDE_HEIGHT_IN = 7.5
MARGIN_LEFT = 0.2
MARGIN_RIGHT = 0.2
GAP = 0.16
ROW_TITLE_HEIGHT = 0.44
TITLE_TOP = 0.16
TITLE_LEFT = 0.26
TITLE_WIDTH = 12.7
IMAGE_TOP = 1.05
IMAGE_BOX_HEIGHT = 5.2
FOOTER_TOP = 6.50
FOOTER_HEIGHT = 0.4
DEFAULT_HARD_SCENARIOS = (
    "filters-gauss-01-b",
    "filters-diffuse-01-f",
    "filters-specular-01-f",
    "filters-light-01-f",
    "filters-light-02-f",
    "filters-overview-02-b",
    "filters-overview-03-b",
    "filters-conv-05-f",
    "text-tspan-01-b",
    "text-tspan-02-b",
    "text-text-07-t",
    "text-text-09-t",
    "coords-trans-09-t",
    "styling-css-01-b",
)


@dataclass
class ScenarioResult:
    name: str
    svg_path: Path
    left_image: Path | None
    right_image: Path | None
    left_status: str
    right_status: str
    left_detail: str | None = None
    right_detail: str | None = None


def _parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument(
        "scenarios",
        nargs="*",
        help="Optional W3C scenario names. Defaults to a curated hard subset.",
    )
    parser.add_argument(
        "--all",
        action="store_true",
        help="Include all tests from tests/svg instead of the hard subset.",
    )
    parser.add_argument(
        "--limit",
        type=int,
        default=None,
        help="Optional limit for the selected scenario list.",
    )
    parser.add_argument(
        "--output",
        type=Path,
        default=Path("reports/visual/w3c-side-by-side"),
        help="Output directory for deck and artifacts.",
    )
    parser.add_argument(
        "--name",
        default="w3c-side-by-side.pptx",
        help="Output deck filename.",
    )
    parser.add_argument(
        "--renderer",
        choices=("soffice", "powerpoint"),
        default="soffice",
        help="Renderer to produce png captures.",
    )
    parser.add_argument(
        "--soffice",
        help="Explicit path to soffice binary (defaults to PATH lookup).",
    )
    parser.add_argument(
        "--verbose",
        action="store_true",
        help="Enable debug logging.",
    )
    parser.add_argument(
        "--template",
        type=Path,
        default=None,
        help="Optional PowerPoint template to inherit design/branding from.",
    )
    parser.add_argument(
        "--ignore-stale-references",
        action="store_true",
        help="Use W3C PNG oracles even when flagged as stale.",
    )
    parser.add_argument(
        "--reference-mode",
        choices=("oracle", "powerpoint"),
        default="oracle",
        help="Left-side reference baseline: W3C PNG oracle or PowerPoint svgBlip capture.",
    )
    parser.add_argument(
        "--force",
        action="store_true",
        help=(
            "Rebuild PPTX/PPTX captures even when existing cache artifacts exist "
            "for the same scenario."
        ),
    )
    parser.add_argument(
        "--no-margins",
        action="store_true",
        help="Use full-slide side-by-side layout with no title/caption margins.",
    )
    return parser.parse_args()


def _require_python_pptx() -> None:
    if Presentation is None:
        raise RuntimeError(
            "python-pptx is required for W3C side-by-side deck generation."
        ) from _PPTX_IMPORT_ERROR


def _discover_w3c_paths() -> dict[str, Path]:
    return {
        path.stem: path
        for path in sorted(Path("tests/svg").glob("*.svg"))
        if path.is_file()
    }


def _build_placeholder(path: Path, title: str, detail: str, size: tuple[int, int]) -> Path:
    image = Image.new("RGB", size, color=(245, 245, 248))
    draw = ImageDraw.Draw(image)
    font = ImageFont.load_default()
    draw.rectangle((16, 16, size[0] - 16, size[1] - 16), outline=(180, 80, 80), width=4)
    draw.text((30, 28), title, font=font, fill=(120, 20, 20))
    wrapped: list[str] = []
    chunk = ""
    for token in detail.split():
        candidate = token if not chunk else f"{chunk} {token}"
        if len(candidate) <= 80:
            chunk = candidate
            continue
        wrapped.append(chunk)
        chunk = token
    if chunk:
        wrapped.append(chunk)
    for row, line in enumerate(wrapped[:20], start=0):
        draw.text((30, 72 + row * 18), line, font=font, fill=(70, 70, 70))
    image.save(path)
    image.close()
    return path


def _is_cached_powerpoint_reference_valid(
    *,
    scenario_paths: list[Path],
    presentation_path: Path,
    render_dir: Path,
    scenario_count: int,
) -> bool:
    if not presentation_path.is_file():
        return False
    try:
        source_mtime = max(path.stat().st_mtime for path in scenario_paths)
    except OSError:
        return False
    if presentation_path.stat().st_mtime < source_mtime:
        return False
    for index in range(scenario_count):
        screenshot = render_dir / f"slide_{index + 1}.png"
        if not screenshot.is_file():
            return False
        if screenshot.stat().st_mtime < presentation_path.stat().st_mtime:
            return False
    return True


def _build_powerpoint_reference_map(
    scenario_items: list[tuple[str, Path]],
    artifacts_root: Path,
    *,
    force_rebuild: bool,
) -> dict[str, Path]:
    if not scenario_items:
        return {}

    reference_root = artifacts_root / "powerpoint"
    reference_root.mkdir(parents=True, exist_ok=True)
    reference_pptx = reference_root / "powerpoint-reference.pptx"
    reference_render_dir = reference_root / "render"
    reference_render_dir.mkdir(parents=True, exist_ok=True)

    source_paths = [path for _name, path in scenario_items]
    scenario_count = len(scenario_items)
    scenario_names = [name for name, _ in scenario_items]

    if not force_rebuild and _is_cached_powerpoint_reference_valid(
        scenario_paths=source_paths,
        presentation_path=reference_pptx,
        render_dir=reference_render_dir,
        scenario_count=scenario_count,
    ):
        return {
            name: reference_render_dir / f"slide_{index + 1}.png"
            for index, name in enumerate(scenario_names)
        }

    for legacy in reference_render_dir.glob("*.png"):
        try:
            legacy.unlink()
        except OSError:
            pass

    embed_svg_collection(source_paths, reference_pptx)
    ref_renderer = resolve_renderer(renderer_name="powerpoint")
    if not ref_renderer.available:
        raise RuntimeError(
            "PowerPoint renderer is not available. This requires PowerPoint + Automation.",
        )

    rendered = ref_renderer.render(reference_pptx, reference_render_dir)
    reference_images = sorted(rendered.images)
    if len(reference_images) < scenario_count:
        raise RuntimeError(
            "PowerPoint reference capture returned fewer images than scenarios."
        )

    return {
        name: reference_images[index]
        for index, name in enumerate(scenario_names)
        if index < len(reference_images)
    }


def _is_cached_render_valid(
    *,
    source_svg: Path,
    pptx_path: Path,
    render_image: Path,
) -> bool:
    if not render_image.is_file() or not pptx_path.is_file():
        return False
    try:
        source_mtime = source_svg.stat().st_mtime
    except OSError:
        return False
    return (
        render_image.stat().st_mtime >= source_mtime
        and pptx_path.stat().st_mtime >= source_mtime
    )


def _discover_render_artifact(
    render_dir: Path,
    scenario_name: str,
) -> Path | None:
    preferred = render_dir / f"{scenario_name}.png"
    if preferred.is_file():
        return preferred
    legacy = render_dir / "presentation.png"
    if legacy.is_file():
        return legacy
    first = next(iter(sorted(render_dir.glob("*.png"))), None)
    return first


def _fit_dimensions(image_path: Path, *, max_width_in: float, max_height_in: float) -> tuple[float, float]:
    with Image.open(image_path) as image:
        width_px = max(1, image.width)
        height_px = max(1, image.height)
    scale = min(max_width_in / width_px, max_height_in / height_px)
    return width_px * scale, height_px * scale


def _add_textbox(
    slide,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    font_size: int,
    bold: bool = False,
    color: tuple[int, int, int] = (0, 0, 0),
    align: str | None = None,
) -> None:
    shape = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    paragraph = shape.text_frame.paragraphs[0]
    paragraph.text = text
    if align == "center":
        paragraph.alignment = PP_ALIGN.CENTER
    elif align == "right":
        paragraph.alignment = PP_ALIGN.RIGHT
    run = paragraph.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)


def _add_fitted_picture(
    slide,
    image_path: Path,
    *,
    left_in: float,
    top_in: float,
    box_width_in: float,
    box_height_in: float,
) -> None:
    width_in, height_in = _fit_dimensions(
        image_path,
        max_width_in=box_width_in,
        max_height_in=box_height_in,
    )
    offset_x = left_in + ((box_width_in - width_in) / 2.0)
    offset_y = top_in + ((box_height_in - height_in) / 2.0)
    slide.shapes.add_picture(
        str(image_path),
        Inches(offset_x),
        Inches(offset_y),
        width=Inches(width_in),
        height=Inches(height_in),
    )


def _render_scenario(
    renderer: PptxRenderer,
    builder: PptxBuilder,
    scenario_name: str,
    svg_path: Path,
    artifacts_root: Path,
    reference_mode: str,
    reference_image: Path | None,
    *,
    ignore_stale_references: bool,
    force_rebuild: bool,
) -> ScenarioResult:
    working_dir = artifacts_root / scenario_name
    working_dir.mkdir(parents=True, exist_ok=True)

    if reference_mode == "powerpoint":
        if reference_image is None:
            left_status = "reference-missing"
            left_detail = "No PowerPoint reference image found for this scenario."
            left_image = _build_placeholder(
                working_dir / "reference.png",
                title="PowerPoint reference missing",
                detail=left_detail,
                size=(1240, 700),
            )
        else:
            left_image = reference_image
            left_status = "powerpoint-reference"
            left_detail = "PowerPoint svgBlip render"
    else:
        stale_reason = _stale_w3c_reference_reason(svg_path)
        reference_path = None
        if ignore_stale_references and stale_reason:
            candidate = svg_path.parent.parent / "png" / f"{svg_path.stem}.png"
            if candidate.is_file():
                reference_path = candidate
        else:
            reference_path = _w3c_reference_png_for_svg(svg_path)

        if reference_path is None:
            reason = stale_reason
            if reason:
                left_status = "missing-reference"
                left_detail = reason
            else:
                left_status = "missing-reference"
                left_detail = "No local W3C PNG oracle found for this fixture."
            left_image = _build_placeholder(
                working_dir / "reference.png",
                title="W3C PNG missing",
                detail=left_detail,
                size=(1240, 700),
            )
        else:
            left_image = working_dir / "reference.png"
            left_image.write_bytes(reference_path.read_bytes())
            left_status = "reference-ok"
            left_detail = f"Using stale reference ({stale_reason})" if stale_reason else None

    render_dir = working_dir / "render"
    render_dir.mkdir(exist_ok=True)
    pptx_path = working_dir / f"{scenario_name}.pptx"
    render_artifact = _discover_render_artifact(render_dir, scenario_name)
    if render_artifact is None:
        render_artifact = render_dir / "presentation.png"
    if not force_rebuild and _is_cached_render_valid(
        source_svg=svg_path, pptx_path=pptx_path, render_image=render_artifact
    ):
        right_image = render_artifact
        right_status = "svg2ooxml"
        right_detail = "cached"
    else:
        try:
            svg_text = svg_path.read_text(encoding="utf-8")
            builder.build_from_svg(svg_text, pptx_path, source_path=svg_path)
            rendered = renderer.render(pptx_path, render_dir)
            if not rendered.images:
                raise FileNotFoundError("No rendered slide images returned.")
            right_image = render_dir / Path(rendered.images[0]).name
            if not right_image.exists():
                right_image = next(iter(render_dir.glob("*.png")))
            right_status = "svg2ooxml"
            right_detail = None
        except Exception as exc:  # pragma: no cover - integration edge
            right_image = _build_placeholder(
                working_dir / "render-error.png",
                title="svg2ooxml render failed",
                detail=str(exc),
                size=(1240, 700),
            )
            right_status = "render-error"
            right_detail = str(exc)

    return ScenarioResult(
        name=scenario_name,
        svg_path=svg_path,
        left_image=left_image,
        right_image=right_image,
        left_status=left_status,
        right_status=right_status,
        left_detail=left_detail,
        right_detail=right_detail,
    )


def _resolve_scenarios(
    args: argparse.Namespace,
    available: dict[str, Path],
) -> list[tuple[str, Path]]:
    if args.all:
        names = [path.stem for path in sorted(Path("tests/svg").glob("*.svg"))]
    else:
        names = [*DEFAULT_HARD_SCENARIOS] if not args.scenarios else list(args.scenarios)

    unknown: list[str] = [name for name in names if name not in available]
    if unknown:
        raise SystemExit(f"Unknown scenario(s): {', '.join(unknown)}")

    unique: list[str] = []
    for name in names:
        if name not in unique:
            unique.append(name)

    if args.limit is not None:
        if args.limit < 0:
            raise SystemExit("--limit must be >= 0")
        unique = unique[: args.limit]

    return [(name, available[name]) for name in unique]


def _pick_template_blank_layout(presentation: Presentation) -> object:
    for layout in presentation.slide_layouts:
        if (layout.name or "").strip().lower() == "blank":
            return layout
    for layout in presentation.slide_layouts:
        if not layout.placeholders:
            return layout
    return presentation.slide_layouts[0]


def _clear_template_slides(presentation: Presentation) -> int:
    """Remove existing template slides and their relationships before adding new ones."""
    removed = 0
    for slide_id in list(presentation.slides._sldIdLst):
        rel_id = getattr(slide_id, "rId", None)
        if rel_id and rel_id in presentation.part._rels:
            presentation.part.drop_rel(rel_id)
        presentation.slides._sldIdLst.remove(slide_id)
        removed += 1
    return removed


def _build_deck(
    results: list[ScenarioResult],
    *,
    output_dir: Path,
    deck_name: str,
    base_presentation: Presentation | None = None,
    no_margins: bool = False,
    reference_mode: str = "oracle",
) -> Path:
    _require_python_pptx()

    presentation = base_presentation if base_presentation is not None else Presentation()
    template_removed = 0
    if base_presentation is not None:
        template_removed = _clear_template_slides(presentation)
        if template_removed:
            logger.info(
                "Removed %s existing template slide(s) before building deck.",
                template_removed,
            )
    if base_presentation is not None:
        blank = _pick_template_blank_layout(presentation)
    elif len(presentation.slide_layouts) > 6:
        blank = presentation.slide_layouts[6]
    else:
        blank = _pick_template_blank_layout(presentation)

    if blank is None:
        raise RuntimeError("No slide layout available for deck generation.")

    presentation.slide_width = Inches(SLIDE_WIDTH_IN)
    presentation.slide_height = Inches(SLIDE_HEIGHT_IN)

    if not no_margins:
        left_label = (
            "Left: W3C reference PNG" if reference_mode == "oracle" else "Left: PowerPoint svgBlip"
        )
        title_slide = presentation.slides.add_slide(blank)
        _add_textbox(
            title_slide,
            left=TITLE_LEFT,
            top=TITLE_TOP,
            width=12.2,
            height=ROW_TITLE_HEIGHT,
            text="W3C Side-by-Side Deck",
            font_size=32,
            bold=True,
        )
        _add_textbox(
            title_slide,
            left=TITLE_LEFT,
            top=0.70,
            width=12.2,
            height=0.35,
            text=f"{left_label}; Right: svg2ooxml PPTX render",
            font_size=14,
            color=(86, 86, 86),
        )
        _add_textbox(
            title_slide,
            left=TITLE_LEFT,
            top=1.06,
            width=12.2,
            height=0.3,
            text=f"Scenarios: {len(results)}",
            font_size=12,
            color=(86, 86, 86),
        )

    if no_margins:
        column_width = SLIDE_WIDTH_IN / 2
        image_top = 0.0
        image_box_height = SLIDE_HEIGHT_IN
    else:
        column_width = (SLIDE_WIDTH_IN - MARGIN_LEFT - MARGIN_RIGHT - GAP) / 2
        image_top = IMAGE_TOP
        image_box_height = IMAGE_BOX_HEIGHT

    for result in results:
        slide = presentation.slides.add_slide(blank)

        if not no_margins:
            _add_textbox(
                slide,
                left=0.25,
                top=0.16,
                width=TITLE_WIDTH,
                height=ROW_TITLE_HEIGHT,
                text=result.name,
                font_size=18,
                bold=True,
                align="left",
            )
            _add_textbox(
                slide,
                left=0.25,
                top=0.53,
                width=12.4,
                height=0.35,
                text=str(result.svg_path),
                font_size=10,
                color=(95, 95, 95),
            )

        left = 0.0 if no_margins else MARGIN_LEFT
        right = left + column_width + (0.0 if no_margins else GAP)

        _add_fitted_picture(
            slide,
            result.left_image,
            left_in=left,
            top_in=image_top,
            box_width_in=column_width,
            box_height_in=image_box_height,
        )
        _add_fitted_picture(
            slide,
            result.right_image,
            left_in=right,
            top_in=image_top,
            box_width_in=column_width,
            box_height_in=image_box_height,
        )

        if not no_margins:
            left_footer = f"{result.left_status}"
            if result.left_detail:
                left_footer = f"{left_footer}: {result.left_detail}"
            right_footer = f"{result.right_status}"
            if result.right_detail:
                right_footer = f"{right_footer}: {result.right_detail}"

            _add_textbox(
                slide,
                left=left,
                top=FOOTER_TOP,
                width=column_width,
                height=FOOTER_HEIGHT,
                text=left_footer,
                font_size=9,
                color=(90, 90, 90),
                align="center",
            )
            _add_textbox(
                slide,
                left=right,
                top=FOOTER_TOP,
                width=column_width,
                height=FOOTER_HEIGHT,
                text=right_footer,
                font_size=9,
                color=(90, 90, 90),
                align="center",
            )

    output_dir.mkdir(parents=True, exist_ok=True)
    deck_path = output_dir / deck_name
    presentation.save(deck_path)
    return deck_path


def main() -> None:
    args = _parse_args()
    logging.basicConfig(
        level=logging.DEBUG if args.verbose else logging.INFO,
        format="%(levelname)s %(message)s",
    )

    _require_python_pptx()

    available = _discover_w3c_paths()
    scenario_items = _resolve_scenarios(args, available)
    if not scenario_items:
        raise SystemExit("No scenarios selected.")

    reference_images: dict[str, Path] = {}
    if args.reference_mode == "powerpoint":
        reference_images = _build_powerpoint_reference_map(
            scenario_items,
            args.output / "artifacts",
            force_rebuild=args.force,
        )

    renderer = resolve_renderer(renderer_name=args.renderer, soffice_path=args.soffice)
    if not renderer.available:
        raise SystemExit(f"Renderer {args.renderer!r} is not available on this system.")

    builder = PptxBuilder(
        filter_strategy="resvg",
        geometry_mode="resvg",
        slide_size_mode="same",
    )

    results: list[ScenarioResult] = []
    for name, path in scenario_items:
        logger.info("Rendering scenario %s", name)
        results.append(
            _render_scenario(
                renderer,
                builder,
                name,
                path,
                args.output / "artifacts",
                args.reference_mode,
                reference_images.get(name),
                ignore_stale_references=args.ignore_stale_references,
                force_rebuild=args.force,
            )
        )

    base_presentation = None
    if args.template is not None:
        if not args.template.is_file():
            raise SystemExit(f"Template not found: {args.template}")
        base_presentation = Presentation(str(args.template))
        logger.info("Using template: %s", args.template)

    deck_path = _build_deck(
        results,
        output_dir=args.output,
        deck_name=args.name,
        base_presentation=base_presentation,
        no_margins=args.no_margins,
        reference_mode=args.reference_mode,
    )
    logger.info("Wrote deck: %s", deck_path)


if __name__ == "__main__":
    main()
